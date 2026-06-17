#!/usr/bin/env python3
"""Generate a professional product introduction PPTX for 需求匹配 (CollabMatch)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ——— Color Palette ———
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xFA, 0xFA, 0xFC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_DARK = RGBColor(0x33, 0x33, 0x44)
GRAY = RGBColor(0x66, 0x66, 0x77)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0xAA)
GRAY_BG = RGBColor(0xF0, 0xF0, 0xF5)
GRAY_CARD = RGBColor(0xF5, 0xF5, 0xFA)
BLUE = RGBColor(0x4F, 0x46, 0xE5)
BLUE_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
BLUE_DARK = RGBColor(0x37, 0x33, 0xBF)
GREEN = RGBColor(0x05, 0x9C, 0x69)
YELLOW = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)

def add_rect(slide, left, top, width, height, fill_color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, width, height, text, size=14, bold=False, color=GRAY_DARK, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box

def add_bullets(slide, left, top, width, height, items, size=12, color=GRAY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return box

def add_card(slide, left, top, width, height, items, title_size=14):
    """Light card with border"""
    card = add_rect(slide, left, top, width, height, GRAY_CARD, 0.05)
    # border
    sp = card._element
    spPr = sp.find(qn('p:spPr'))
    ln = spPr.makeelement(qn('a:ln'), {})
    ln.set('w', '6350')  # 0.5pt
    srgb = ln.makeelement(qn('a:srgbClr'), {'val': 'E5E7EB'})
    ln.append(srgb)
    spPr.append(ln)
    
    top_y = top + Inches(0.2)
    for item in items:
        typ, text = item[0], item[1]
        if typ == 'title':
            add_text(slide, left + Inches(0.2), top_y, width - Inches(0.4), Inches(0.4),
                     text, size=title_size, bold=True, color=DARK)
            top_y += Inches(0.4)
        elif typ == 'desc':
            add_text(slide, left + Inches(0.2), top_y, width - Inches(0.4), Inches(0.5),
                     text, size=10, color=GRAY)
            top_y += Inches(0.4)
        elif typ == 'icon':
            add_text(slide, left + Inches(0.2), top_y, width - Inches(0.4), Inches(0.5),
                     text, size=22)
            top_y += Inches(0.35)
        elif typ == 'bullets':
            add_bullets(slide, left + Inches(0.2), top_y, width - Inches(0.4), height - (top_y - top),
                       text, size=10, color=GRAY)
    return card


# ================================================================
# SLIDE 1 — Cover
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
# Dark header block
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(4.2), DARK)
# Accent bar
add_rect(sl, Inches(0), Inches(4.2), Inches(13.33), Pt(4), BLUE)

add_text(sl, Inches(1), Inches(0.6), Inches(11), Inches(0.4),
         '产品介绍', size=14, color=GRAY_LIGHT)
add_text(sl, Inches(1), Inches(1.3), Inches(11), Inches(1),
         '需求匹配', size=54, bold=True, color=WHITE)
add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
         'AI 驱动的协作匹配平台', size=20, color=RGBColor(0xAA, 0xAA, 0xCC))
add_text(sl, Inches(1), Inches(3.3), Inches(11), Inches(0.5),
         '说一个想法，找到志同道合的伙伴 —— 从灵感碰撞到团队组建，一站式完成',
         size=13, color=GRAY_LIGHT)

# Info tags
tags = ['AI 智能匹配', '多领域协作', '即时沟通', '游戏化体验', '零框架依赖']
for i, t in enumerate(tags):
    x = Inches(1) + Inches(i * 2.3)
    add_rect(sl, x, Inches(5.0), Inches(2.0), Inches(0.45), BLUE_LIGHT, 0.3)
    add_text(sl, x + Inches(0.15), Inches(5.02), Inches(1.7), Inches(0.4),
             t, size=11, color=BLUE, align=PP_ALIGN.CENTER)

add_text(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
         '2026 年 6 月  ·  产品团队', size=12, color=GRAY_LIGHT)

# ================================================================
# SLIDE 2 — The Problem
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '个人与项目的协作困境', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

problems = [
    ('🧑‍💻  独立开发者', '有技术能力但缺少设计、运营等互补伙伴，一个人扛所有角色', '技术侧'),
    ('🎨  设计师/创作者', '有创意但找不到靠谱的技术实现者，好想法止步于概念', '创意侧'),
    ('🎓  学生/创业团队', '组队靠熟人，圈子有限，很难找到真正匹配的跨界伙伴', '校园侧'),
    ('💼  自由职业者', '项目不稳定，缺少持续的项目来源和稳定的协作网络', '自由侧'),
]
for i, (title, desc, tag) in enumerate(problems):
    y = Inches(1.5) + Inches(i * 1.45)
    add_rect(sl, Inches(1), y, Inches(11.33), Inches(1.25), GRAY_CARD, 0.05)
    add_text(sl, Inches(1.3), y + Inches(0.15), Inches(3), Inches(0.35),
             title, size=15, bold=True, color=DARK)
    add_text(sl, Inches(1.3), y + Inches(0.55), Inches(8), Inches(0.45),
             desc, size=11, color=GRAY)
    add_rect(sl, Inches(10.5), y + Inches(0.15), Inches(1.5), Inches(0.3), BLUE_LIGHT, 0.3)
    add_text(sl, Inches(10.5), y + Inches(0.17), Inches(1.5), Inches(0.3),
             tag, size=9, color=BLUE, align=PP_ALIGN.CENTER)

# Highlight
add_rect(sl, Inches(1), Inches(7.0), Inches(11.33), Inches(0.4), RGBColor(0xFE, 0xF2, 0xF2), 0.05)
add_text(sl, Inches(1.2), Inches(7.0), Inches(11), Inches(0.4),
         '💡  核心痛点：协作连接成本高，没有高效的工具来发现和匹配对的人',
         size=12, bold=True, color=RED)

# ================================================================
# SLIDE 3 — Vision
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '产品愿景', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         '让每个人都能轻松找到志同道合的伙伴，让每个好想法都能找到落地的团队',
         size=18, color=GRAY_DARK)

visions = [
    ('🎯', '精准匹配', 'AI 深入分析需求与技能互补度，推荐最合适的伙伴，提高协作成功率'),
    ('🚀', '低门槛启动', 'AI 辅助生成结构化需求文档，一键发布，无需从零开始'),
    ('🎮', '游戏化体验', '像素风驿站场景，让匹配像玩游戏一样自然有趣'),
    ('🔄', '双向连接', '既可发布需求找人，也可浏览项目找机会，双向匹配覆盖更多场景'),
]
for i, (icon, title, desc) in enumerate(visions):
    x = Inches(0.7) + Inches(i * 3.1)
    add_rect(sl, x, Inches(2.6), Inches(2.9), Inches(3.2), GRAY_CARD, 0.05)
    add_text(sl, x + Inches(0.2), Inches(2.8), Inches(2.5), Inches(0.4),
             icon, size=28)
    add_text(sl, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(0.3),
             title, size=15, bold=True, color=DARK)
    add_text(sl, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(1.8),
             desc, size=11, color=GRAY)

# ================================================================
# SLIDE 4 — Core Flow
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '核心流程：从想法到团队', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

steps = ['💡 描述想法', '🤖 AI 解析', '🔍 智能匹配', '💬 沟通协作', '🤝 组队成功']
for i, s in enumerate(steps):
    x = Inches(0.4) + Inches(i * 2.6)
    add_rect(sl, x, Inches(1.6), Inches(2.3), Inches(1.2), BLUE_LIGHT, 0.3)
    add_text(sl, x + Inches(0.1), Inches(1.7), Inches(2.1), Inches(0.4),
             s, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(sl, x + Inches(2.3), Inches(1.85), Inches(0.3), Inches(0.4),
                 '→', size=20, color=BLUE)

features = [
    ('📋  AI 辅助需求整理', '输入你的想法，AI 自动生成结构化需求文档，包含项目背景、目标、所需技能、时间线等'),
    ('🎯  双向智能匹配', '既可发布需求找人，也可浏览项目找机会。AI 分析技能互补度和协作意愿，推荐最合适的伙伴'),
    ('🔄  一键组队协作', '匹配确认后自动创建协作群，内置即时通讯，支持文件分享、AI 群聊总结、成员管理'),
]
for i, (title, desc) in enumerate(features):
    y = Inches(3.1) + Inches(i * 1.35)
    add_rect(sl, Inches(0.7), y, Inches(11.93), Inches(1.15), GRAY_CARD, 0.05)
    add_text(sl, Inches(1), y + Inches(0.1), Inches(4), Inches(0.3),
             title, size=14, bold=True, color=DARK)
    add_text(sl, Inches(1), y + Inches(0.5), Inches(11), Inches(0.5),
             desc, size=11, color=GRAY)

# ================================================================
# SLIDE 5 — Key Features
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '核心功能矩阵', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

features = [
    ('🤖', 'AI 智能匹配', '正向找人 + 反向找项目\n双向匹配，AI 推荐'),
    ('💬', 'AI 对话', '内置 LLM 生成需求\n支持诊断与优化'),
    ('📊', '需求广场', '按领域浏览公开需求\n发现感兴趣的项目'),
    ('👥', '协作群组', '自动建群，实时聊天\nAI 总结，文件分享'),
    ('👤', '个人主页', '展示技能、作品集、\n协作历史与评价'),
    ('🎮', '游戏化驿站', '2D 像素风场景\nNPC 交互触发匹配'),
    ('🔧', 'AI 工具箱', '生成 PRD / SWOT 分析\n周期估算等 7+ 工具'),
    ('📡', '协作者入口', 'MCP 协议接入\nAI Agent 自动协作'),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(col * 3.15)
    y = Inches(1.4) + Inches(row * 3.0)
    add_rect(sl, x, y, Inches(2.95), Inches(2.7), GRAY_CARD, 0.05)
    add_text(sl, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.4),
             icon, size=26)
    add_text(sl, x + Inches(0.2), y + Inches(0.6), Inches(2.5), Inches(0.3),
             title, size=14, bold=True, color=DARK)
    add_text(sl, x + Inches(0.2), y + Inches(1.0), Inches(2.5), Inches(1.2),
             desc, size=10, color=GRAY)

# ================================================================
# SLIDE 6 — Game Innovation
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '游戏化创新：「驿站」', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.4),
         '把协作匹配变成一场像素风 RPG 探索体验', size=16, color=GRAY_DARK)

games = [
    ('🏠', '驿站场景', '20×15 Tile 的古风驿站，包含接待台、公告栏、茶座、比武擂台等区域，每个区域对应不同的功能入口'),
    ('🧑‍🦱', 'NPC 交互', '与"小二""侠客"等 NPC 对话触发匹配面板，游戏化替代冰冷表单'),
    ('🎵', '沉浸体验', 'Web Audio API 合成脚步声、环境音效、交互反馈音效，增强沉浸感'),
    ('🔄', '零外部依赖', '自研 Canvas 游戏引擎，6 个模块，纯原生 Web API，无需额外加载'),
]
for i, (icon, title, desc) in enumerate(games):
    x = Inches(0.5) + Inches(i * 3.2)
    add_rect(sl, x, Inches(2.2), Inches(3.0), Inches(2.5), GRAY_CARD, 0.05)
    add_text(sl, x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(0.4),
             icon, size=32)
    add_text(sl, x + Inches(0.2), Inches(2.9), Inches(2.6), Inches(0.3),
             title, size=15, bold=True, color=DARK)
    add_text(sl, x + Inches(0.2), Inches(3.3), Inches(2.6), Inches(1.2),
             desc, size=10, color=GRAY)

# Core concept highlight
add_rect(sl, Inches(1), Inches(5.0), Inches(11.33), Inches(1.5), BLUE_LIGHT, 0.05)
add_text(sl, Inches(1.3), Inches(5.15), Inches(10.7), Inches(0.3),
         '✨ 核心理念', size=16, bold=True, color=BLUE)
add_text(sl, Inches(1.3), Inches(5.55), Inches(10.7), Inches(0.8),
         '看似网络游戏，实则联通现实 —— 游戏内的每一次交互（对话 NPC、发布需求、组队匹配）都对应真实的平台操作，'
         '通过游戏化的方式降低用户心理门槛，让协作匹配变得有趣且自然。',
         size=12, color=GRAY_DARK)

# ================================================================
# SLIDE 7 — Tech Stack
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '技术架构', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

# Frontend
add_rect(sl, Inches(0.7), Inches(1.5), Inches(5.8), Inches(3.0), GRAY_CARD, 0.05)
add_text(sl, Inches(1), Inches(1.6), Inches(5.2), Inches(0.3),
         '🌐  前端', size=16, bold=True, color=DARK)
add_bullets(sl, Inches(1), Inches(2.1), Inches(5.2), Inches(2.2),
           ['单文件 SPA，零框架依赖，极速加载',
            'Lucide SVG 图标系统，视觉风格统一',
            '中英双语国际化，彩蛋式江湖模式切换',
            '自研 2D Canvas 游戏引擎（6 个模块）',
            '零外部依赖，纯原生 Web API',
            'GameBus 事件总线连接游戏与 DOM'],
           size=11, color=GRAY)

# Backend
add_rect(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.0), GRAY_CARD, 0.05)
add_text(sl, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.3),
         '⚙️  后端与基础设施', size=16, bold=True, color=DARK)
add_bullets(sl, Inches(7.1), Inches(2.1), Inches(5.2), Inches(2.2),
           ['腾讯云 CloudBase 云函数，弹性伸缩',
            '火山方舟（豆包）LLM API，AI 对话',
            'CloudBase 内置数据库（MongoDB 兼容）',
            'HTTP 访问服务统一路由',
            '双平台部署：CloudBase + Vercel',
            'MongoDB Atlas 香港节点（待迁移）'],
           size=11, color=GRAY)

# Tags
tags = ['Vite', 'Lucide', 'CloudBase', '豆包 API', 'Canvas', 'Web Audio', 'MongoDB Atlas', 'JWT', 'Vercel']
for i, t in enumerate(tags):
    x = Inches(0.7) + Inches(i * 1.38)
    if x + Inches(1.2) > Inches(12.6):
        break
    add_rect(sl, x, Inches(4.8), Inches(1.25), Inches(0.4), BLUE_LIGHT, 0.3)
    add_text(sl, x + Inches(0.05), Inches(4.82), Inches(1.15), Inches(0.35),
             t, size=9, color=BLUE, align=PP_ALIGN.CENTER)

# Architecture note
add_rect(sl, Inches(0.7), Inches(5.5), Inches(11.93), Inches(1.5), GRAY_CARD, 0.05)
add_text(sl, Inches(1), Inches(5.6), Inches(11.3), Inches(0.3),
         '📐  架构特点', size=14, bold=True, color=DARK)
add_bullets(sl, Inches(1), Inches(6.0), Inches(11.3), Inches(1.0),
           ['• 前后端分离：静态托管 + 云函数，可独立扩缩',
            '• 全栈 Serverless：无需管理服务器，按量计费',
            '• 双平台部署：国内 CloudBase，国际 Vercel，覆盖更多用户',
            '• 低耦合设计：数据库层可平滑迁移（内置 DB ↔ MongoDB Atlas）'],
           size=11, color=GRAY)

# ================================================================
# SLIDE 8 — Roadmap
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '规划路线图', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

roadmap = [
    ('🚀  MVP 已完成', '邮箱注册 · 需求发布 · AI 对话 · 智能匹配 · 协作群组 · 游戏驿站 · 中英双语', GREEN, '已完成'),
    ('🔧  正式上线进行中', '绑定自定义域名 · 用户协议/隐私政策 · 内容审核 · 国际化部署', BLUE, '进行中'),
    ('📋  短期规划', '短信验证码登录 · 邮箱验证 · 密码找回 · 用户实名 · 作品集上传', YELLOW, '规划中'),
    ('🎯  中期规划', '迁移 MongoDB Atlas · WebSocket 实时通信 · 移动端适配 · 成就徽章系统', RGBColor(0x25, 0x63, 0xEB), '规划中'),
    ('🌟  长期规划', '更多游戏场景（茶馆·擂台） · 智能合约结算 · AI Agent 撮合 · 开放平台 API', RED, '规划中'),
]
for i, (title, desc, color, status) in enumerate(roadmap):
    y = Inches(1.5) + Inches(i * 1.15)
    # Status dot
    add_rect(sl, Inches(1), y + Inches(0.1), Pt(10), Pt(10), color)
    add_text(sl, Inches(1.4), y, Inches(3.5), Inches(0.35),
             title, size=15, bold=True, color=DARK)
    add_text(sl, Inches(4.8), y, Inches(7), Inches(0.35),
             desc, size=11, color=GRAY)
    # Status badge
    badge_color = BLUE_LIGHT if status == '进行中' else GRAY_CARD
    add_rect(sl, Inches(11.5), y + Inches(0.03), Inches(0.8), Inches(0.28), badge_color, 0.3)
    add_text(sl, Inches(11.5), y + Inches(0.04), Inches(0.8), Inches(0.26),
             status, size=8, color=GRAY, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 9 — Business Model
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.0), DARK)
add_text(sl, Inches(1), Inches(0.25), Inches(11), Inches(0.5),
         '商业模式探讨', size=24, bold=True, color=WHITE)
add_rect(sl, Inches(1), Inches(1.2), Inches(11.33), Pt(2), BLUE)

biz = [
    ('🏅', '个人增值服务', '基础匹配免费，增值按需付费',
     ['金牌协作者认证标识', '每月 AI 匹配次数扩容', '作品集展示优化', '协作评价与信用体系']),
    ('🏢', '企业/团队版', '面向组织内部的协作匹配需求',
     ['企业内部协作匹配系统', '项目管理看板集成', '私有化部署支持', 'SSO / 权限管理']),
    ('🔌', 'AI Agent 生态', '面向 AI Agent 的协作市场',
     ['MCP 协议接入服务费', 'AI Agent 自动发布/应标', 'Agent 间自动撮合交易', '平台小额抽佣']),
    ('🌐', '社区与流量', '面向广泛用户的基础商业',
     ['项目展示推荐位', '付费置顶需求', '协作者推荐排名', '增值 SaaS 工具订阅']),
]
for i, (icon, title, subtitle, items) in enumerate(biz):
    x = Inches(0.4) + Inches(i * 3.2)
    add_rect(sl, x, Inches(1.5), Inches(3.0), Inches(4.5), GRAY_CARD, 0.05)
    add_text(sl, x + Inches(0.2), Inches(1.7), Inches(2.6), Inches(0.4),
             icon, size=28)
    add_text(sl, x + Inches(0.2), Inches(2.15), Inches(2.6), Inches(0.3),
             title, size=15, bold=True, color=DARK)
    add_text(sl, x + Inches(0.2), Inches(2.5), Inches(2.6), Inches(0.3),
             subtitle, size=10, color=BLUE)
    add_bullets(sl, x + Inches(0.2), Inches(2.9), Inches(2.6), Inches(1.5),
               [f'• {item}' for item in items], size=10, color=GRAY)

add_rect(sl, Inches(1), Inches(6.3), Inches(11.33), Inches(0.8), BLUE_LIGHT, 0.05)
add_text(sl, Inches(1.3), Inches(6.4), Inches(10.7), Inches(0.5),
         '💡  核心理念：基础匹配免费，增值服务付费 —— 降低门槛，靠价值变现',
         size=14, bold=True, color=BLUE)

# ================================================================
# SLIDE 10 — Q&A
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(7.5), DARK)

add_text(sl, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
         '感谢聆听', size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(6), Inches(2.7), Inches(1.33), Pt(3), BLUE)
add_text(sl, Inches(1), Inches(3.0), Inches(11), Inches(0.5),
         '期待你的反馈和建议', size=20, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

contacts = [
    ('📱  产品演示', 'CloudBase 在线体验'),
    ('📧  联系我们', '产品团队'),
    ('💻  GitHub', 'github.com/whongchen9/collabmatch'),
]
for i, (name, url) in enumerate(contacts):
    x = Inches(1.5) + Inches(i * 3.5)
    add_rect(sl, x, Inches(4.5), Inches(3.0), Inches(1.2), RGBColor(0x22, 0x22, 0x3A), 0.05)
    add_text(sl, x + Inches(0.2), Inches(4.7), Inches(2.6), Inches(0.3),
             name, size=14, bold=True, color=WHITE)
    add_text(sl, x + Inches(0.2), Inches(5.1), Inches(2.6), Inches(0.3),
             url, size=11, color=GRAY_LIGHT)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '需求匹配-产品介绍.pptx')
prs.save(output_path)
print(f'✅ PPTX saved to: {output_path}')
