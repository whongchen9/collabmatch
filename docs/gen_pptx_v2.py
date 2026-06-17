#!/usr/bin/env python3
"""Generate a clean, minimal PPTX that renders well on both desktop and mobile"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARK = RGBColor(0x1A,0x1A,0x2E)
GRAY = RGBColor(0x66,0x66,0x77)
LIGHT = RGBColor(0x99,0x99,0xAA)
BLUE = RGBColor(0x4F,0x46,0xE5)
HEADER_BG = RGBColor(0x1A,0x1A,0x2E)

def tb(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Simple text box helper"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def header(slide, title):
    """Dark header bar"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_BG
    shape.line.fill.background()
    tb(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.5), title, size=22, bold=True, color=WHITE)

def body_text(slide, text, top=1.2, size=14, color=GRAY, bold=False):
    tb(slide, Inches(0.8), Inches(top), Inches(11.7), Inches(0.5), text, size=size, bold=bold, color=color)

def dark_slide(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = HEADER_BG

# ===== SLIDE 1: Cover =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(sl)
tb(sl, Inches(0.8), Inches(2.2), Inches(11), Inches(0.4), '产品介绍', size=14, color=LIGHT, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(3.0), Inches(11), Inches(1.0), '需求匹配', size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5), 'AI 驱动的协作匹配平台', size=20, color=RGBColor(0xAA,0xAA,0xCC), align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5), '说一个想法，找到志同道合的伙伴', size=14, color=LIGHT, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4), '2026 年 6 月  ·  产品团队', size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: Problem =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '个人与项目的协作困境')
problems = [
    ('独立开发者', '有技术能力但缺少设计、运营等互补伙伴，一个人扛所有角色'),
    ('设计师 / 创作者', '有创意但找不到靠谱的技术实现者，好想法止步于概念'),
    ('学生 / 创业团队', '组队靠熟人，圈子有限，很难找到真正匹配的跨界伙伴'),
    ('自由职业者', '项目不稳定，缺少持续的项目来源和稳定的协作网络'),
]
for i, (t, d) in enumerate(problems):
    y = 1.2 + i * 1.2
    tb(sl, Inches(0.8), Inches(y), Inches(3.5), Inches(0.4), t, size=16, bold=True, color=DARK)
    tb(sl, Inches(0.8), Inches(y+0.4), Inches(11), Inches(0.4), d, size=12, color=GRAY)
tb(sl, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4), '💡 核心痛点：协作连接成本高，缺少高效的工具来发现和匹配对的人', size=13, bold=True, color=RGBColor(0xDC,0x26,0x26))

# ===== SLIDE 3: Vision =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '产品愿景')
tb(sl, Inches(0.8), Inches(1.3), Inches(11), Inches(0.8), '让每个人都能轻松找到志同道合的伙伴，让每个好想法都能找到落地的团队', size=18, color=DARK)
visions = [('🎯 精准匹配', 'AI 分析需求与技能互补度，推荐最合适的伙伴'),
           ('🚀 低门槛启动', 'AI 辅助生成需求文档，无需从零开始'),
           ('🎮 游戏化体验', '像素风驿站场景，让匹配像玩游戏一样自然'),
           ('🔄 双向连接', '发布需求找人，也可浏览项目找机会')]
for i, (t, d) in enumerate(visions):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.5 + row * 1.8
    tb(sl, Inches(x), Inches(y), Inches(5.5), Inches(0.35), t, size=15, bold=True, color=DARK)
    tb(sl, Inches(x), Inches(y+0.4), Inches(5.5), Inches(0.4), d, size=12, color=GRAY)

# ===== SLIDE 4: Flow =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '核心流程：从想法到团队')
steps = ['💡 描述想法', '→', '🤖 AI 解析', '→', '🔍 智能匹配', '→', '💬 沟通协作', '→', '🤝 组队成功']
line = ''
for s in steps:
    line += s + '  '
tb(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5), line, size=16, bold=True, color=BLUE)
features = [
    ('📋 AI 辅助需求整理', '输入想法，AI 自动生成结构化需求文档，包含项目背景、目标、所需技能、时间线等'),
    ('🎯 双向智能匹配', '找人找项目双向匹配，分析技能互补度和协作意愿，推荐最合适的伙伴'),
    ('🔄 一键组队协作', '匹配确认后自动创建协作群，内置即时通讯，支持文件分享和 AI 群聊总结'),
]
for i, (t, d) in enumerate(features):
    y = 2.5 + i * 1.5
    tb(sl, Inches(0.8), Inches(y), Inches(11), Inches(0.35), t, size=15, bold=True, color=DARK)
    tb(sl, Inches(0.8), Inches(y+0.4), Inches(11), Inches(0.5), d, size=12, color=GRAY)

# ===== SLIDE 5: Features =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '核心功能矩阵')
feats = [
    ('🤖', 'AI 智能匹配', '正向找人+反向找项目'),
    ('💬', 'AI 对话', 'LLM 生成需求文档'),
    ('📊', '需求广场', '按领域浏览公开需求'),
    ('👥', '协作群组', '自动建群，实时聊天'),
    ('👤', '个人主页', '技能、作品集展示'),
    ('🎮', '游戏化驿站', '像素风 NPC 交互'),
    ('🔧', 'AI 工具箱', 'PRD/SWOT 等 7+ 工具'),
    ('📡', '协作者入口', 'MCP 协议接入'),
]
for i, (icon, t, d) in enumerate(feats):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.15
    y = 1.3 + row * 2.8
    tb(sl, Inches(x), Inches(y), Inches(2.8), Inches(0.35), icon + '  ' + t, size=15, bold=True, color=DARK)
    tb(sl, Inches(x), Inches(y+0.4), Inches(2.8), Inches(0.4), d, size=11, color=GRAY)

# ===== SLIDE 6: Game =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '游戏化创新：「驿站」')
tb(sl, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4), '把协作匹配变成一场像素风 RPG 探索体验', size=15, color=GRAY)
games = [
    ('🏠 驿站场景', '20×15 Tile 的古风驿站：接待台、公告栏、茶座、比武擂台'),
    ('🧑‍🦱 NPC 交互', '与"小二""侠客"对话触发匹配面板，游戏化替代冰冷表单'),
    ('🎵 沉浸体验', 'Web Audio API 合成脚步声、环境音效、交互音效'),
    ('🔄 零外部依赖', '自研 Canvas 游戏引擎，纯原生 Web API，无需额外加载'),
]
for i, (t, d) in enumerate(games):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 2.0 + row * 1.6
    tb(sl, Inches(x), Inches(y), Inches(5.5), Inches(0.35), t, size=15, bold=True, color=DARK)
    tb(sl, Inches(x), Inches(y+0.4), Inches(5.5), Inches(0.5), d, size=12, color=GRAY)
tb(sl, Inches(0.8), Inches(5.5), Inches(11), Inches(1.0),
   '💡 核心理念：看似网络游戏，实则联通现实 —— 游戏内的每一次交互都对应真实的平台操作，让协作匹配变得有趣且自然。',
   size=13, bold=True, color=BLUE)

# ===== SLIDE 7: Tech =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '技术架构')
tb(sl, Inches(0.8), Inches(1.3), Inches(3), Inches(0.35), '🌐 前端', size=16, bold=True, color=DARK)
tb(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.5),
   '• 单文件 SPA，零框架依赖\n• Lucide SVG 图标系统，视觉统一\n• 中英双语国际化，彩蛋式江湖模式切换\n• 自研 2D Canvas 游戏引擎（6 个模块）\n• 零外部依赖，纯原生 Web API',
   size=12, color=GRAY)
tb(sl, Inches(7), Inches(1.3), Inches(3), Inches(0.35), '⚙️ 后端与基础设施', size=16, bold=True, color=DARK)
tb(sl, Inches(7), Inches(1.8), Inches(5.5), Inches(1.5),
   '• 腾讯云 CloudBase 云函数，弹性伸缩\n• 火山方舟（豆包）LLM API，AI 对话\n• CloudBase 内置数据库（MongoDB 兼容）\n• HTTP 访问服务统一路由\n• 双平台部署：CloudBase + Vercel',
   size=12, color=GRAY)
tb(sl, Inches(0.8), Inches(3.8), Inches(11), Inches(0.3),
   '技术标签：Vite  ·  Lucide  ·  CloudBase  ·  豆包 API  ·  Canvas  ·  Web Audio  ·  JWT  ·  Vercel',
   size=11, color=BLUE)
tb(sl, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
   '📐 架构特点：\n• 前后端分离：静态托管 + 云函数，可独立扩缩\n• 全栈 Serverless：无需管理服务器，按量计费\n• 双平台覆盖：国内 CloudBase，国际 Vercel\n• 低耦合设计：数据库层可平滑迁移（内置 DB ↔ MongoDB Atlas）',
   size=12, color=GRAY)

# ===== SLIDE 8: Roadmap =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '规划路线图')
roadmap = [
    ('🚀  MVP 已完成', '邮箱注册 · 需求发布 · AI 对话 · 智能匹配 · 协作群组 · 游戏驿站 · 中英双语', '✅ 已完成'),
    ('🔧  正式上线进行中', '绑定自定义域名 · 用户协议/隐私政策 · 内容审核 · 国际化部署', '🟡 进行中'),
    ('📋  短期规划', '短信验证码登录 · 邮箱验证 · 密码找回 · 用户实名 · 作品集上传', '📍 规划中'),
    ('🎯  中期规划', '迁移 MongoDB Atlas · WebSocket 实时通信 · 移动端适配 · 成就徽章系统', '📍 规划中'),
    ('🌟  长期规划', '更多游戏场景（茶馆·擂台） · 智能合约结算 · AI Agent 撮合 · 开放平台 API', '📍 规划中'),
]
for i, (t, d, s) in enumerate(roadmap):
    y = 1.2 + i * 1.15
    tb(sl, Inches(0.8), Inches(y), Inches(4), Inches(0.35), t, size=15, bold=True, color=DARK)
    tb(sl, Inches(0.8), Inches(y+0.35), Inches(10), Inches(0.35), d, size=11, color=GRAY)
    tb(sl, Inches(11.5), Inches(y), Inches(1.5), Inches(0.35), s, size=10, color=GRAY)

# ===== SLIDE 9: Business =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, '商业模式探讨')
biz = [
    ('🏅 个人增值服务', '基础匹配免费，增值按需付费。金牌认证 / 匹配扩容 / 作品集优化 / 信用体系'),
    ('🏢 企业/团队版', '面向组织内部。企业内部匹配 / 项目看板 / 私有部署 / SSO 权限管理'),
    ('🔌 AI Agent 生态', '面向 AI Agent 的协作市场。MCP 接入费 / Agent 自动撮合 / 平台微抽佣'),
    ('🌐 社区与流量', '面向广泛用户。推荐位广告 / 付费置顶 / 推荐排名 / SaaS 工具订阅'),
]
for i, (t, d) in enumerate(biz):
    y = 1.3 + i * 1.4
    tb(sl, Inches(0.8), Inches(y), Inches(5), Inches(0.35), t, size=16, bold=True, color=DARK)
    tb(sl, Inches(0.8), Inches(y+0.4), Inches(11), Inches(0.5), d, size=12, color=GRAY)
tb(sl, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
   '💡 核心理念：基础匹配免费，增值服务付费 —— 降低门槛，靠价值变现', size=14, bold=True, color=BLUE)

# ===== SLIDE 10: QA =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_slide(sl)
tb(sl, Inches(0.8), Inches(2.5), Inches(11), Inches(0.5), '感谢聆听', size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5), '期待你的反馈和建议', size=20, color=LIGHT, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4), '产品演示 · CloudBase 部署 · GitHub: github.com/whongchen9/collabmatch', size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), '需求匹配-产品介绍.pptx')
prs.save(out)
print(f'✅ 已保存: {out}')
